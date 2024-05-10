
from pptx import Presentation
from pptx.util import Inches, Pt

def slideShow(fileName):
    pr = Presentation()
    return pr

def savePP(pr, fileName):
    pr.save(fileName)

def getTopicSlide(pr, topic, infoTypes):
    slide_register = pr.slide_layouts[6]
    slide = pr.slides.add_slide(slide_register)

    left = top = Inches(.25)
    width = height = Inches(0)
    txBox = slide.shapes.add_textbox(left, top, width, height)

    topicName = str(topic[0]) + " " + str(topic[1])
    addTopicTBX(txBox, topicName)
    
    infoBox = txBox.text_frame
    numTypes = len(infoTypes)

    index = 2
    while index < numTypes:
        addInfoLine(infoBox, infoTypes[index], topic[index])
        index += 1


def addTopicTBX(txBox, topicName):
    topicBX = txBox.text_frame
    p = topicBX.add_paragraph()
    p.text = str(topicName)
    p.font.size = Pt(30)


def addInfoLine(infoBox, infoType, info):
    p = infoBox.add_paragraph()
    p.text = str(infoType) + ": " + str(info)
    p.font.size = Pt(20)
    p.level = 1
